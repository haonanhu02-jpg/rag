"""Independent format parsers that normalize all input into one block contract."""

from __future__ import annotations

import mimetypes
from io import BytesIO
from itertools import groupby
from typing import Any
from uuid import UUID, uuid5

import pdfplumber
import pypdfium2 as pdfium
from bs4 import BeautifulSoup, Tag
from charset_normalizer import from_bytes
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from markdown_it import MarkdownIt
from openpyxl import load_workbook
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from rag_platform.adapters.outbound.document_security import validate_image_size, validate_ooxml
from rag_platform.adapters.outbound.ocr import ocr_lines
from rag_platform.domain.identifiers import BlockId
from rag_platform.modules.knowledge.contracts import (
    BinaryDocumentParser,
    BlockKind,
    BoundingBox,
    CompiledBlock,
    CoordinateSpace,
    DocumentParseError,
    DocumentResourceLimit,
    MediaReference,
    OcrEngine,
    ParsedPayload,
    ParserRequest,
    ParseWarning,
    TableMetadata,
)

_BLOCK_NAMESPACE = UUID("291638d2-22af-53bb-a0fe-994d15f33d38")


class _Blocks:
    def __init__(self, request: ParserRequest, parser_name: str, parser_version: str) -> None:
        self._request = request
        self._parser_name = parser_name
        self._parser_version = parser_version
        self._values: list[CompiledBlock] = []
        self._cursor = 0
        self._headings: list[str] = []

    def add(
        self,
        kind: BlockKind,
        text: str = "",
        *,
        page: int | None = None,
        box: BoundingBox | None = None,
        table: TableMetadata | None = None,
        media: MediaReference | None = None,
        confidence: float | None = None,
        heading_level: int | None = None,
        warnings: tuple[ParseWarning, ...] = (),
    ) -> None:
        value = text.strip()
        if kind != BlockKind.IMAGE and not value:
            return
        if kind == BlockKind.HEADING:
            level = max(1, heading_level or 1)
            self._headings = self._headings[: level - 1]
            self._headings.append(value)
        heading_path = tuple(self._headings)
        ordinal = len(self._values)
        start, end = self._cursor, self._cursor + len(value)
        identity = uuid5(
            _BLOCK_NAMESPACE,
            "\x1f".join(
                (
                    str(self._request.document_id),
                    self._request.source_sha256,
                    self._parser_name,
                    self._parser_version,
                    str(ordinal),
                    str(kind),
                    str(page),
                    value,
                    media.embedded_path if media and media.embedded_path else "",
                )
            ),
        )
        self._values.append(
            CompiledBlock(
                id=BlockId(identity),
                ordinal=ordinal,
                kind=kind,
                text=value,
                start_character=start,
                end_character=end,
                page_number=page,
                bounding_box=box,
                heading_path=heading_path,
                table=table,
                media=media,
                confidence=confidence,
                parser_name=self._parser_name,
                parser_version=self._parser_version,
                warnings=warnings,
            )
        )
        self._cursor = end + 2

    def done(
        self, warnings: tuple[ParseWarning, ...] = (), page_count: int | None = None
    ) -> ParsedPayload:
        return ParsedPayload(
            self._parser_name,
            self._parser_version,
            tuple(self._values),
            warnings,
            page_count,
        )


class TextDocumentParser(BinaryDocumentParser):
    format_ids = frozenset({"text", "markdown"})
    parser_name = "structured-text"
    parser_version = "2"

    def parse(self, content: bytes, request: ParserRequest) -> ParsedPayload:
        decoded = from_bytes(content).best()
        if decoded is None or decoded.percent_chaos > 30:
            raise DocumentParseError("parser_text_encoding", "text encoding is not recognizable")
        text = str(decoded).replace("\r\n", "\n").replace("\r", "\n").strip()
        if not text:
            raise DocumentParseError("parser_no_content", "text document is empty")
        blocks = _Blocks(request, self.parser_name, self.parser_version)
        if request.format_id == "text":
            for value in text.split("\n\n"):
                blocks.add(BlockKind.PARAGRAPH, value)
            return blocks.done()
        tokens = MarkdownIt("commonmark").enable("table").parse(text)
        index = 0
        while index < len(tokens):
            token = tokens[index]
            if token.type == "heading_open" and index + 1 < len(tokens):
                blocks.add(
                    BlockKind.HEADING,
                    tokens[index + 1].content,
                    heading_level=int(token.tag[1:]),
                )
                index += 3
                continue
            if token.type == "paragraph_open" and index + 1 < len(tokens):
                blocks.add(BlockKind.PARAGRAPH, tokens[index + 1].content)
                index += 3
                continue
            if token.type in {"fence", "code_block"}:
                blocks.add(BlockKind.CODE, token.content)
            elif token.type == "list_item_open":
                values: list[str] = []
                cursor = index + 1
                depth = 1
                while cursor < len(tokens) and depth:
                    if tokens[cursor].type == "list_item_open":
                        depth += 1
                    elif tokens[cursor].type == "list_item_close":
                        depth -= 1
                    elif tokens[cursor].type == "inline" and depth == 1:
                        values.append(tokens[cursor].content)
                    cursor += 1
                blocks.add(BlockKind.LIST, "\n".join(values))
                index = cursor
                continue
            elif token.type == "table_open":
                rows: list[list[str]] = []
                row: list[str] = []
                cursor = index + 1
                while cursor < len(tokens) and tokens[cursor].type != "table_close":
                    current = tokens[cursor]
                    if current.type == "tr_open":
                        row = []
                    elif current.type == "inline":
                        row.append(current.content)
                    elif current.type == "tr_close" and row:
                        rows.append(row)
                    cursor += 1
                self._add_table(blocks, rows)
                index = cursor + 1
                continue
            index += 1
        return blocks.done()

    @staticmethod
    def _add_table(blocks: _Blocks, rows: list[list[str]]) -> None:
        if not rows:
            return
        columns = max(len(row) for row in rows)
        blocks.add(
            BlockKind.TABLE,
            "\n".join("\t".join(row) for row in rows),
            table=TableMetadata(len(rows), columns, True),
        )


class HtmlDocumentParser(BinaryDocumentParser):
    format_ids = frozenset({"html"})
    parser_name = "safe-html"
    parser_version = "2"

    def parse(self, content: bytes, request: ParserRequest) -> ParsedPayload:
        decoded = from_bytes(content).best()
        if decoded is None:
            raise DocumentParseError("parser_html_encoding", "HTML encoding is not recognizable")
        soup = BeautifulSoup(str(decoded), "html.parser")
        for active in soup.find_all(("script", "style", "noscript", "template")):
            active.decompose()
        blocks = _Blocks(request, self.parser_name, self.parser_version)
        selected = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "table", "img"}
        root = soup.body or soup
        for tag in root.find_all(selected):
            if any(parent.name in selected for parent in tag.parents if isinstance(parent, Tag)):
                continue
            name = tag.name
            if name.startswith("h"):
                blocks.add(
                    BlockKind.HEADING,
                    tag.get_text(" ", strip=True),
                    heading_level=int(name[1]),
                )
            elif name == "p":
                blocks.add(BlockKind.PARAGRAPH, tag.get_text(" ", strip=True))
            elif name == "li":
                blocks.add(BlockKind.LIST, tag.get_text(" ", strip=True))
            elif name == "pre":
                blocks.add(BlockKind.CODE, tag.get_text("\n", strip=True))
            elif name == "table":
                rows = [
                    [cell.get_text(" ", strip=True) for cell in row.find_all(("th", "td"))]
                    for row in tag.find_all("tr")
                ]
                TextDocumentParser._add_table(blocks, [row for row in rows if row])
            elif name == "img":
                path = str(tag.get("src") or "embedded-image")
                media_type = mimetypes.guess_type(path)[0] or "application/octet-stream"
                blocks.add(
                    BlockKind.IMAGE,
                    str(tag.get("alt") or ""),
                    media=MediaReference(media_type, path),
                )
        return blocks.done()


class DocxDocumentParser(BinaryDocumentParser):
    format_ids = frozenset({"docx"})
    parser_name = "structured-docx"
    parser_version = "2"

    def parse(self, content: bytes, request: ParserRequest) -> ParsedPayload:
        validate_ooxml(content, request.limits)
        try:
            document = Document(BytesIO(content))
        except Exception as exc:
            raise DocumentParseError("parser_docx_invalid", "invalid DOCX document") from exc
        blocks = _Blocks(request, self.parser_name, self.parser_version)
        for child in document.element.body.iterchildren():
            if isinstance(child, CT_P):
                paragraph = Paragraph(child, document)
                text = paragraph.text.strip()
                style = (paragraph.style.name if paragraph.style else "").casefold()
                if text:
                    if style.startswith("heading"):
                        suffix = style.rsplit(" ", maxsplit=1)[-1]
                        blocks.add(
                            BlockKind.HEADING,
                            text,
                            heading_level=int(suffix) if suffix.isdigit() else 1,
                        )
                    elif "list" in style or (
                        paragraph._p.pPr is not None and paragraph._p.pPr.numPr is not None
                    ):
                        blocks.add(BlockKind.LIST, text)
                    elif "code" in style:
                        blocks.add(BlockKind.CODE, text)
                    else:
                        blocks.add(BlockKind.PARAGRAPH, text)
                for blip in paragraph._p.xpath(".//a:blip"):
                    relationship_id = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    part = document.part.related_parts.get(relationship_id)
                    if part is None:
                        continue
                    path = str(part.partname).lstrip("/")
                    blocks.add(
                        BlockKind.IMAGE,
                        media=MediaReference(part.content_type, path),
                    )
            elif isinstance(child, CT_Tbl):
                table = Table(child, document)
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                TextDocumentParser._add_table(blocks, rows)
        return blocks.done()


class PptxDocumentParser(BinaryDocumentParser):
    format_ids = frozenset({"pptx"})
    parser_name = "structured-pptx"
    parser_version = "2"

    def parse(self, content: bytes, request: ParserRequest) -> ParsedPayload:
        validate_ooxml(content, request.limits)
        try:
            presentation = Presentation(BytesIO(content))
        except Exception as exc:
            raise DocumentParseError("parser_pptx_invalid", "invalid PPTX document") from exc
        if len(presentation.slides) > request.limits.max_pages:
            raise DocumentResourceLimit("pptx_slides", "presentation has too many slides")
        blocks = _Blocks(request, self.parser_name, self.parser_version)
        width = float(presentation.slide_width or 1)
        height = float(presentation.slide_height or 1)
        for page, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                box = BoundingBox(
                    max(0.0, float(shape.left) / width),
                    max(0.0, float(shape.top) / height),
                    min(1.0, float(shape.left + shape.width) / width),
                    min(1.0, float(shape.top + shape.height) / height),
                    CoordinateSpace.NORMALIZED,
                )
                if shape.has_table:
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    if rows:
                        blocks.add(
                            BlockKind.TABLE,
                            "\n".join("\t".join(row) for row in rows),
                            page=page,
                            box=box,
                            table=TableMetadata(len(rows), max(map(len, rows)), True),
                        )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    blocks.add(
                        BlockKind.IMAGE,
                        page=page,
                        box=box,
                        media=MediaReference(
                            image.content_type,
                            image.filename,
                            image.size[0],
                            image.size[1],
                        ),
                    )
                elif shape.has_text_frame and shape.text.strip():
                    kind = BlockKind.HEADING if shape == slide.shapes.title else BlockKind.PARAGRAPH
                    blocks.add(kind, shape.text, page=page, box=box, heading_level=1)
        return blocks.done(page_count=len(presentation.slides))


class XlsxDocumentParser(BinaryDocumentParser):
    format_ids = frozenset({"xlsx"})
    parser_name = "structured-xlsx"
    parser_version = "2"

    def parse(self, content: bytes, request: ParserRequest) -> ParsedPayload:
        validate_ooxml(content, request.limits)
        try:
            workbook = load_workbook(
                BytesIO(content), read_only=False, data_only=False, keep_links=False
            )
        except Exception as exc:
            raise DocumentParseError("parser_xlsx_invalid", "invalid XLSX workbook") from exc
        try:
            if len(workbook.worksheets) > request.limits.max_worksheets:
                raise DocumentResourceLimit("xlsx_sheets", "workbook has too many worksheets")
            blocks = _Blocks(request, self.parser_name, self.parser_version)
            warnings: list[ParseWarning] = []
            formula = False
            merged = False
            total_cells = 0
            for page, sheet in enumerate(workbook.worksheets, start=1):
                total_cells += sheet.max_row * sheet.max_column
                if total_cells > request.limits.max_spreadsheet_cells:
                    raise DocumentResourceLimit("xlsx_cells", "workbook has too many cells")
                rows: list[list[str]] = []
                for values in sheet.iter_rows(values_only=True):
                    row = ["" if value is None else str(value) for value in values]
                    while row and not row[-1]:
                        row.pop()
                    if row:
                        formula = formula or any(value.startswith("=") for value in row)
                        rows.append(row)
                merged = merged or bool(sheet.merged_cells.ranges)
                if rows:
                    columns = max(len(row) for row in rows)
                    blocks.add(
                        BlockKind.TABLE,
                        f"Sheet: {sheet.title}\n" + "\n".join("\t".join(row) for row in rows),
                        page=page,
                        table=TableMetadata(len(rows), columns, True),
                    )
            if formula:
                warnings.append(
                    ParseWarning(
                        "xlsx_formula_preserved_not_evaluated",
                        "formulas are preserved as source text and never evaluated",
                    )
                )
            if merged:
                warnings.append(
                    ParseWarning("xlsx_merged_cells_present", "merged cells were flattened")
                )
            return blocks.done(tuple(warnings), len(workbook.worksheets))
        finally:
            workbook.close()


class PdfDocumentParser(BinaryDocumentParser):
    format_ids = frozenset({"pdf"})
    parser_name = "layout-pdf"
    parser_version = "2"

    def __init__(self, ocr: OcrEngine) -> None:
        self._ocr = ocr

    def parse(self, content: bytes, request: ParserRequest) -> ParsedPayload:
        try:
            document = pdfplumber.open(BytesIO(content))
        except Exception as exc:
            raise DocumentParseError("parser_pdf_invalid", "invalid PDF document") from exc
        warnings: list[ParseWarning] = []
        blocks = _Blocks(request, self.parser_name, self.parser_version)
        try:
            if len(document.pages) > request.limits.max_pages:
                raise DocumentResourceLimit("pdf_pages", "PDF has too many pages")
            for page_number, page in enumerate(document.pages, start=1):
                words = page.extract_words(x_tolerance=2, y_tolerance=3, keep_blank_chars=False)
                if words:
                    self._add_pdf_words(blocks, words, page_number)
                else:
                    image = self._render(content, page_number - 1)
                    validate_image_size(*image.size, request.limits)
                    result = self._ocr.recognize(
                        image,
                        language=request.ocr_language,
                        timeout_seconds=request.limits.ocr_timeout_seconds,
                    )
                    lines = ocr_lines(result)
                    if not lines:
                        raise DocumentParseError("ocr_no_text", "OCR produced no text")
                    for text, box, confidence in lines:
                        blocks.add(
                            BlockKind.PARAGRAPH,
                            text,
                            page=page_number,
                            box=box,
                            confidence=confidence,
                        )
                    warnings.append(
                        ParseWarning(
                            "pdf_scanned_page_ocr",
                            "page had no text layer and used OCR",
                            page_number,
                        )
                    )
                for table in page.find_tables():
                    rows = [
                        ["" if cell is None else str(cell).strip() for cell in row]
                        for row in table.extract()
                    ]
                    if not rows:
                        continue
                    x0, top, x1, bottom = map(float, table.bbox)
                    blocks.add(
                        BlockKind.TABLE,
                        "\n".join("\t".join(row) for row in rows),
                        page=page_number,
                        box=BoundingBox(x0, top, x1, bottom, CoordinateSpace.PAGE_POINTS),
                        table=TableMetadata(len(rows), max(map(len, rows)), True),
                    )
            return blocks.done(tuple(warnings), len(document.pages))
        finally:
            document.close()

    @staticmethod
    def _add_pdf_words(blocks: _Blocks, words: list[dict[str, Any]], page: int) -> None:
        ordered = sorted(words, key=lambda word: (round(float(word["top"]), 1), float(word["x0"])))
        for _, line_values in groupby(ordered, key=lambda word: round(float(word["top"]) / 3) * 3):
            line = list(line_values)
            segment: list[dict[str, Any]] = []
            segments: list[list[dict[str, Any]]] = []
            for word in line:
                if segment and float(word["x0"]) - float(segment[-1]["x1"]) > 80:
                    segments.append(segment)
                    segment = []
                segment.append(word)
            if segment:
                segments.append(segment)
            for values in segments:
                blocks.add(
                    BlockKind.PARAGRAPH,
                    " ".join(str(word["text"]) for word in values),
                    page=page,
                    box=BoundingBox(
                        min(float(word["x0"]) for word in values),
                        min(float(word["top"]) for word in values),
                        max(float(word["x1"]) for word in values),
                        max(float(word["bottom"]) for word in values),
                        CoordinateSpace.PAGE_POINTS,
                    ),
                )

    @staticmethod
    def _render(content: bytes, page: int) -> Image.Image:
        try:
            document = pdfium.PdfDocument(content)
            rendered: Image.Image = document[page].render(scale=2).to_pil()
            rendered.load()
            document.close()
            return rendered
        except Exception as exc:
            raise DocumentParseError("parser_pdf_render", "PDF page could not be rendered") from exc


class ImageDocumentParser(BinaryDocumentParser):
    format_ids = frozenset({"image"})
    parser_name = "image-ocr"
    parser_version = "2"

    def __init__(self, ocr: OcrEngine) -> None:
        self._ocr = ocr

    def parse(self, content: bytes, request: ParserRequest) -> ParsedPayload:
        try:
            image = Image.open(BytesIO(content))
            validate_image_size(*image.size, request.limits)
            image.load()
        except (UnidentifiedImageError, OSError) as exc:
            raise DocumentParseError("parser_image_invalid", "invalid image") from exc
        blocks = _Blocks(request, self.parser_name, self.parser_version)
        width, height = image.size
        blocks.add(
            BlockKind.IMAGE,
            media=MediaReference(
                image.get_format_mimetype() or request.media_type,
                request.file_name,
                width,
                height,
            ),
            page=1,
            box=BoundingBox(0, 0, float(width), float(height), CoordinateSpace.PIXELS),
        )
        result = self._ocr.recognize(
            image,
            language=request.ocr_language,
            timeout_seconds=request.limits.ocr_timeout_seconds,
        )
        lines = ocr_lines(result)
        if not lines:
            raise DocumentParseError("ocr_no_text", "OCR produced no text")
        for text, box, confidence in lines:
            blocks.add(
                BlockKind.PARAGRAPH,
                text,
                page=1,
                box=box,
                confidence=confidence,
            )
        return blocks.done(page_count=1)


def build_document_parsers(ocr: OcrEngine) -> tuple[BinaryDocumentParser, ...]:
    return (
        TextDocumentParser(),
        HtmlDocumentParser(),
        DocxDocumentParser(),
        PptxDocumentParser(),
        XlsxDocumentParser(),
        PdfDocumentParser(ocr),
        ImageDocumentParser(ocr),
    )
