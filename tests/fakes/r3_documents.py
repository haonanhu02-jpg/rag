"""Generated legal format samples and deterministic OCR for R3 tests."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas

from rag_platform.modules.knowledge.contracts import (
    BoundingBox,
    CoordinateSpace,
    OcrResult,
    OcrWord,
)

DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


@dataclass(frozen=True, slots=True)
class FormatSample:
    name: str
    media_type: str
    content: bytes
    expected_kinds: frozenset[str]


class StaticOcr:
    def __init__(self, text: str = "alarm reset inspection") -> None:
        self.text = text
        self.calls: list[tuple[str, float]] = []

    def recognize(
        self, image: object, *, language: str, timeout_seconds: float
    ) -> OcrResult:
        del image
        self.calls.append((language, timeout_seconds))
        return OcrResult(
            "static-ocr",
            "1",
            language,
            tuple(
                OcrWord(
                    word,
                    0.95,
                    index,
                    BoundingBox(
                        float(10 + index * 80),
                        12,
                        float(70 + index * 80),
                        36,
                        CoordinateSpace.PIXELS,
                    ),
                )
                for index, word in enumerate(self.text.split())
            ),
        )

    def available_languages(self) -> frozenset[str]:
        return frozenset({"eng", "chi_sim"})


def samples() -> tuple[FormatSample, ...]:
    return (
        FormatSample(
            "manual.txt",
            "text/plain",
            b"Alarm recovery\n\nReset the controller.",
            frozenset({"paragraph"}),
        ),
        FormatSample(
            "manual.md",
            "text/markdown",
            (
                b"# Alarm Recovery\n\nReset controller.\n\n- Inspect relay\n\n"
                b"| Step | Action |\n| --- | --- |\n| 1 | Isolate power |\n"
            ),
            frozenset({"heading", "paragraph", "list", "table"}),
        ),
        FormatSample(
            "manual.html",
            "text/html",
            (
                b"<h1>Alarm Recovery</h1><p>Reset controller.</p>"
                b"<script>secret()</script><table><tr><th>Step</th><th>Action</th></tr>"
                b"<tr><td>1</td><td>Inspect relay</td></tr></table>"
                b"<img src='relay.png' alt='Relay diagram'>"
            ),
            frozenset({"heading", "paragraph", "table", "image"}),
        ),
        FormatSample(
            "manual.docx",
            DOCX,
            _docx(),
            frozenset({"heading", "paragraph", "table", "image"}),
        ),
        FormatSample(
            "manual.pptx",
            PPTX,
            _pptx(),
            frozenset({"heading", "paragraph", "table", "image"}),
        ),
        FormatSample("alarms.xlsx", XLSX, _xlsx(), frozenset({"table"})),
        FormatSample("manual.pdf", "application/pdf", _pdf(), frozenset({"paragraph", "table"})),
        FormatSample("alarm.png", "image/png", _image(), frozenset({"image", "paragraph"})),
    )


def blank_pdf() -> bytes:
    stream = BytesIO()
    document = canvas.Canvas(stream, pagesize=(612, 792))
    document.showPage()
    document.save()
    return stream.getvalue()


def _image() -> bytes:
    image = Image.new("RGB", (480, 120), "white")
    ImageDraw.Draw(image).text((20, 40), "ALARM RESET INSPECTION", fill="black")
    stream = BytesIO()
    image.save(stream, "PNG")
    return stream.getvalue()


def _docx() -> bytes:
    document = Document()
    document.add_heading("Alarm Recovery", level=1)
    document.add_paragraph("Reset the controller and inspect the relay.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text, table.cell(0, 1).text = "Step", "Action"
    table.cell(1, 0).text, table.cell(1, 1).text = "1", "Isolate power"
    document.add_picture(BytesIO(_image()), width=Inches(1))
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


def _pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Alarm Recovery"
    slide.placeholders[1].text = "Reset the controller."
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(5), Inches(1)).table
    table.cell(0, 0).text, table.cell(0, 1).text = "Step", "Action"
    table.cell(1, 0).text, table.cell(1, 1).text = "1", "Inspect relay"
    slide.shapes.add_picture(BytesIO(_image()), Inches(6), Inches(4), width=Inches(1))
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Alarms"
    sheet.append(["Code", "Action"])
    sheet.append(["A-100", '=CONCAT("Inspect", " relay")'])
    sheet.merge_cells("A4:B4")
    sheet["A4"] = "Approved"
    workbook.create_sheet("Assets").append(["Asset", "Location"])
    stream = BytesIO()
    workbook.save(stream)
    workbook.close()
    return stream.getvalue()


def _pdf() -> bytes:
    stream = BytesIO()
    document = canvas.Canvas(stream, pagesize=(612, 792))
    document.drawString(72, 740, "Alarm Recovery")
    document.drawString(72, 716, "Reset the controller and inspect the relay.")
    left, top, width, height = 72, 660, 160, 24
    for row in range(3):
        document.line(left, top - row * height, left + 2 * width, top - row * height)
    for column in range(3):
        document.line(left + column * width, top, left + column * width, top - 2 * height)
    document.drawString(left + 4, top - 17, "Step")
    document.drawString(left + width + 4, top - 17, "Action")
    document.drawString(left + 4, top - height - 17, "1")
    document.drawString(left + width + 4, top - height - 17, "Isolate power")
    document.save()
    return stream.getvalue()
